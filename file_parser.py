# -*- coding: utf-8 -*-
"""
文件解析模块 - 支持文本、Word、Excel、图片等多种格式
"""

import os
import re
import sys
import logging
import warnings
from typing import List, Dict, Optional, Tuple
from pathlib import Path

# 设置环境变量（Windows）
if sys.platform == 'win32':
    os.environ['PYTHONIOENCODING'] = 'utf-8'

from config import SUPPORTED_EXTENSIONS, MAX_CONTENT_LENGTH

# 禁用不必要的警告
warnings.filterwarnings('ignore', message='Neither CUDA nor MPS are available')
warnings.filterwarnings('ignore', category=UserWarning)

# 配置logging使用UTF-8编码
import sys
if sys.platform == 'win32':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

logging.basicConfig(
    level=logging.WARNING, 
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(sys.stdout)
    ]
)
logger = logging.getLogger(__name__)


class FileParser:
    """文件解析器"""
    
    def __init__(self):
        self.supported_extensions = SUPPORTED_EXTENSIONS
    
    def parse_file(self, file_path: str) -> Optional[Dict]:
        """
        解析单个文件，返回结构化数据
        
        Args:
            file_path: 文件路径
            
        Returns:
            包含文件信息的字典，包括：
            - file_path: 文件路径
            - file_name: 文件名
            - file_type: 文件类型
            - content: 文本内容
            - metadata: 元数据
        """
        if not os.path.exists(file_path):
            logger.warning(f"文件不存在: {file_path}")
            return None
        
        # 检查文件大小，超过5MB忽略
        file_size = os.path.getsize(file_path)
        max_size = 5 * 1024 * 1024  # 5MB
        if file_size > max_size:
            logger.warning(f"文件过大({file_size/1024/1024:.1f}MB > 5MB)，跳过: {os.path.basename(file_path)}")
            return None
        
        ext = Path(file_path).suffix.lower()
        if ext not in self.supported_extensions:
            logger.warning(f"不支持的文件类型: {ext}")
            return None
        
        # 确保文件名使用UTF-8编码
        try:
            file_name = os.path.basename(file_path)
            # 如果文件名是bytes，解码为str
            if isinstance(file_name, bytes):
                file_name = file_name.decode('utf-8', errors='replace')
        except:
            file_name = os.path.basename(file_path)
        
        try:
            if ext == '.txt':
                content = self._parse_txt(file_path)
            elif ext in ['.doc', '.docx']:
                content = self._parse_word(file_path)
            elif ext in ['.xlsx', '.xls']:
                content = self._parse_excel(file_path)
            elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif']:
                content = self._parse_image(file_path)
            elif ext == '.pdf':
                content = self._parse_pdf(file_path)
            elif ext == '.json':
                content = self._parse_json(file_path)
            elif ext == '.csv':
                content = self._parse_csv(file_path)
            elif ext in ['.md', '.markdown']:
                content = self._parse_markdown(file_path)
            elif ext in ['.xml', '.html', '.htm']:
                content = self._parse_html(file_path)
            elif ext == '.zip':
                content = self._parse_zip(file_path)
            elif ext == '.xmind':
                content = self._parse_xmind(file_path)
            elif ext == '.pptx':
                content = self._parse_pptx(file_path)
            elif ext == '.sql':
                content = self._parse_sql(file_path)
            else:
                content = self._parse_txt(file_path)  # 尝试作为文本解析
            
            if not content or len(content.strip()) < 10:
                return None
            
            problem_info = self._extract_problem_info(file_name, content)
            
            return {
                'file_path': file_path,
                'file_name': file_name,
                'file_type': ext,
                'content': content[:MAX_CONTENT_LENGTH],
                'problem_id': problem_info.get('problem_id', ''),
                'problem_type': problem_info.get('problem_type', ''),
                'metadata': {
                    'size': os.path.getsize(file_path),
                    'modified_time': os.path.getmtime(file_path)
                }
            }
            
        except Exception as e:
            # 只记录文件名，不记录完整路径
            logger.warning(f"解析跳过: {file_name}")
            return None
    
    def _parse_txt(self, file_path: str) -> str:
        """解析文本文件"""
        encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                    return self._clean_content(content)
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"读取文本文件失败: {str(e)}")
                break
        
        return ""
    
    def _parse_sql(self, file_path: str) -> str:
        """解析SQL文件 - 直接读取文本"""
        return self._parse_txt(file_path)
    
    def _parse_word(self, file_path: str) -> str:
        """解析Word文档 - 支持.docx和旧版.doc"""
        file_name = os.path.basename(file_path)
        ext = os.path.splitext(file_path)[1].lower()
        
        # 方法1: 使用python-docx (适用于.docx)
        if ext == '.docx':
            try:
                from docx import Document
                doc = Document(file_path)
                content = []
                
                for para in doc.paragraphs:
                    if para.text.strip():
                        content.append(para.text)
                
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            content.append(' | '.join(row_text))
                
                if content:
                    logger.info(f"Word解析成功(docx): {file_name}")
                    return self._clean_content('\n'.join(content))
            except ImportError:
                logger.debug("python-docx未安装")
            except Exception as e:
                logger.debug(f"python-docx解析失败: {e}")
        
        # 方法2: 使用antiword (Linux/Mac上解析.doc)
        try:
            import subprocess
            result = subprocess.run(['antiword', '-m', 'UTF-8.txt', file_path],
                                    capture_output=True, text=True, timeout=10)
            if result.returncode == 0 and result.stdout and result.stdout.strip():
                logger.info(f"Word解析成功(antiword): {file_name}")
                return self._clean_content(result.stdout)
            elif result.returncode != 0 and result.stderr:
                logger.debug(f"antiword返回错误: {result.stderr.strip()}")
        except (ImportError, FileNotFoundError):
            logger.debug("antiword未安装")
        except Exception as e:
            logger.debug(f"antiword解析失败: {e}")
        
        # 方法3: 使用textract (通用文档解析)
        try:
            import textract
            text = textract.process(file_path).decode('utf-8')
            if text and len(text.strip()) > 10:
                logger.info(f"Word解析成功(textract): {file_name}")
                return self._clean_content(text)
        except ImportError:
            logger.debug("textract未安装")
        except Exception as e:
            logger.debug(f"textract解析失败: {e}")
        
        # 方法4: 使用pywin32 (Windows上通过COM接口)
        try:
            import win32com.client
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(file_path)
            content = doc.Content.Text
            doc.Close()
            word.Quit()
            if content and len(content.strip()) > 10:
                logger.info(f"Word解析成功(pywin32): {file_name}")
                return self._clean_content(content)
        except ImportError:
            logger.debug("pywin32未安装")
        except Exception as e:
            logger.debug(f"pywin32解析失败: {e}")
        
        # 方法5: 使用LibreOffice headless转换（最可靠的.doc解析方案）
        try:
            import subprocess
            import tempfile
            import shutil
            tmp_dir = tempfile.mkdtemp()
            try:
                result = subprocess.run(
                    ['soffice', '--headless', '--convert-to', 'txt:Text',
                     '--outdir', tmp_dir, file_path],
                    capture_output=True, text=True, timeout=30
                )
                if result.returncode == 0:
                    base_name = os.path.splitext(file_name)[0]
                    txt_path = os.path.join(tmp_dir, base_name + '.txt')
                    if os.path.exists(txt_path):
                        with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
                            text = f.read()
                        if text and len(text.strip()) > 10:
                            logger.info(f"Word解析成功(LibreOffice): {file_name}")
                            return self._clean_content(text)
            finally:
                shutil.rmtree(tmp_dir, ignore_errors=True)
        except FileNotFoundError:
            logger.debug("LibreOffice未安装")
        except Exception as e:
            logger.debug(f"LibreOffice解析失败: {e}")

        # 方法6: 使用olefile提取旧版.doc文本（改进版）
        if ext == '.doc':
            try:
                import olefile
                ole = olefile.OleFileIO(file_path)
                text_parts = []
                # 尝试从多个流中提取文本
                for stream_name in ['WordDocument', '1Table', '0Table']:
                    if not ole.exists(stream_name):
                        continue
                    try:
                        data = ole.openstream(stream_name).read()
                        raw = data.decode('utf-16-le', errors='ignore')
                        # 保留可打印字符：ASCII + 中文 + 标点等
                        filtered = re.sub(
                            r'[^\u0020-\u007E\u00A0-\u00FF'
                            r'\u4E00-\u9FFF\u3000-\u303F'
                            r'\uFF00-\uFFEF\u2000-\u206F'
                            r'\uAC00-\uD7AF\n\r\t]+', '', raw)
                        filtered = re.sub(r'\s+', ' ', filtered).strip()
                        if len(filtered) > 10:
                            text_parts.append(filtered)
                    except Exception:
                        continue
                ole.close()
                if text_parts:
                    combined = ' '.join(text_parts)
                    if len(combined.strip()) > 10:
                        logger.info(f"Word解析成功(olefile): {file_name}")
                        return self._clean_content(combined)
            except ImportError:
                logger.debug("olefile未安装")
            except Exception as e:
                logger.debug(f"olefile解析失败: {e}")
        
        # 方法7: 使用aspose-words (商业库，但功能强大)
        try:
            import aspose.words as aw
            doc = aw.Document(file_path)
            content = doc.get_text()
            if content and len(content.strip()) > 10:
                logger.info(f"Word解析成功(aspose): {file_name}")
                return self._clean_content(content)
        except ImportError:
            logger.debug("aspose-words未安装")
        except Exception as e:
            logger.debug(f"aspose解析失败: {e}")
        
        logger.warning(f"Word文档解析失败，无可用解析器: {file_name}")
        return ""
    
    def _parse_excel(self, file_path: str) -> str:
        """解析Excel文件"""
        try:
            import pandas as pd
            
            df = pd.read_excel(file_path, sheet_name=None)
            content = []
            
            for sheet_name, sheet_df in df.items():
                content.append(f"=== 工作表: {sheet_name} ===")
                
                for col in sheet_df.columns:
                    if sheet_df[col].notna().any():
                        values = sheet_df[col].dropna().astype(str).tolist()
                        content.append(f"{col}: {', '.join(values[:10])}")
                
                for idx, row in sheet_df.iterrows():
                    row_data = []
                    for col in sheet_df.columns:
                        if pd.notna(row[col]):
                            row_data.append(str(row[col]))
                    if row_data:
                        content.append(' | '.join(row_data))
            
            return self._clean_content('\n'.join(content))
            
        except ImportError:
            logger.warning("pandas库未安装，尝试使用其他方法")
            return ""
        except Exception as e:
            logger.error(f"解析Excel文件失败: {str(e)}")
            return ""
    
    def _parse_image(self, file_path: str) -> str:
        """解析图片文件（OCR）- 支持多种OCR引擎"""
        file_name = os.path.basename(file_path)
        
        # 检查文件大小，太大的图片跳过OCR
        file_size = os.path.getsize(file_path)
        if file_size > 5 * 1024 * 1024:  # 大于5MB
            logger.warning(f"图片过大({file_size/1024/1024:.1f}MB)，跳过OCR: {file_name}")
            return f"[图片文件: {file_name} ({file_size/1024/1024:.1f}MB)]"
        
        # 首先尝试用PIL打开图片验证
        try:
            from PIL import Image
            img = Image.open(file_path)
            img.verify()  # 验证图片
            # 重新打开（verify后需要重新打开）
            img = Image.open(file_path)
            
            # 获取原始尺寸
            orig_size = img.size
            
            # 转换为RGB（处理RGBA、P模式等）
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # 如果图片太大，缩小尺寸以提高OCR速度
            max_dimension = 1500
            if max(orig_size) > max_dimension:
                ratio = max_dimension / max(orig_size)
                new_size = (int(orig_size[0] * ratio), int(orig_size[1] * ratio))
                img = img.resize(new_size, Image.LANCZOS)
                logger.debug(f"图片已缩放: {orig_size} -> {new_size}")
            
            # 保存为临时文件
            import tempfile
            temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
            img.save(temp_file.name, 'PNG')
            temp_path = temp_file.name
            temp_file.close()
        except Exception as e:
            logger.warning(f"图片验证失败: {file_name}, 错误: {e}")
            return f"[图片文件: {file_name} (无法打开)]"
        
        try:
            # 方法1: 使用RapidOCR (最快，ONNXRuntime)
            try:
                from rapidocr_onnxruntime import RapidOCR
                logger.debug(f"尝试使用RapidOCR: {file_name}")
                
                # 初始化OCR（只初始化一次）
                if not hasattr(self, '_rapid_ocr'):
                    logger.info("正在初始化RapidOCR...")
                    self._rapid_ocr = RapidOCR()
                    logger.info("RapidOCR初始化完成")
                
                result, elapse = self._rapid_ocr(temp_path)
                
                if result and len(result) > 0:
                    texts = []
                    for item in result:
                        if isinstance(item, (list, tuple)) and len(item) >= 2:
                            texts.append(item[1])
                    text = '\n'.join(texts)
                    if text and len(text.strip()) > 3:
                        logger.info(f"RapidOCR识别成功: {file_name}, 识别到{len(texts)}行文本")
                        return self._clean_content(text)
                logger.debug(f"RapidOCR未识别到文本: {file_name}")
            except ImportError:
                logger.debug("RapidOCR未安装")
            except Exception as e:
                logger.debug(f"RapidOCR失败: {e}")
            
            # 方法2: 使用easyocr (备选)
            try:
                import easyocr
                logger.debug(f"尝试使用EasyOCR: {file_name}")
                
                # 初始化reader（使用类变量缓存）
                if not hasattr(self, '_easyocr_reader'):
                    logger.info("正在初始化EasyOCR（首次使用需要下载模型，请耐心等待）...")
                    self._easyocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=False, verbose=False)
                    logger.info("EasyOCR初始化完成")
                
                result = self._easyocr_reader.readtext(temp_path, detail=0, paragraph=True)
                
                if result and len(result) > 0:
                    text = '\n'.join(result)
                    if text and len(text.strip()) > 3:
                        logger.info(f"EasyOCR识别成功: {file_name}, 识别到{len(result)}行文本")
                        return self._clean_content(text)
                logger.debug(f"EasyOCR未识别到文本: {file_name}")
            except ImportError:
                logger.debug("EasyOCR未安装")
            except Exception as e:
                logger.debug(f"EasyOCR失败: {e}")
            
            # 方法3: 使用pytesseract (最后备选)
            try:
                import pytesseract
                from PIL import Image
                
                logger.debug(f"尝试使用Tesseract: {file_name}")
                image = Image.open(temp_path)
                image = image.convert('L')  # 灰度图
                
                text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                if text and len(text.strip()) > 3:
                    logger.info(f"Tesseract OCR识别成功: {file_name}")
                    return self._clean_content(text)
                    
                logger.debug(f"Tesseract OCR未识别到文本: {file_name}")
            except ImportError:
                logger.debug("pytesseract未安装")
            except Exception as e:
                logger.debug(f"Tesseract OCR失败: {e}")
            
            logger.warning(f"图片OCR失败，无可用引擎或无法识别: {file_name}")
            logger.warning(f"  建议安装RapidOCR: pip install rapidocr_onnxruntime")
            return f"[图片文件: {file_name} (OCR识别失败)]"
            
        finally:
            # 清理临时文件
            try:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
            except:
                pass
    
    def _parse_pdf(self, file_path: str) -> str:
        """解析PDF文件"""
        file_name = os.path.basename(file_path)
        
        # 方法1: 使用PyPDF2
        try:
            import PyPDF2
            content = []
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        content.append(text)
            if content:
                logger.info(f"PDF解析成功(PyPDF2): {file_name}")
                return self._clean_content('\n'.join(content))
        except ImportError:
            logger.debug("PyPDF2未安装")
        except Exception as e:
            logger.debug(f"PyPDF2解析失败: {e}")
        
        # 方法2: 使用pdfplumber (效果更好)
        try:
            import pdfplumber
            content = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        content.append(text)
            if content:
                logger.info(f"PDF解析成功(pdfplumber): {file_name}")
                return self._clean_content('\n'.join(content))
        except ImportError:
            logger.debug("pdfplumber未安装")
        except Exception as e:
            logger.debug(f"pdfplumber解析失败: {e}")
        
        # 方法3: 使用pdf2image + OCR (扫描版PDF)
        try:
            from pdf2image import convert_from_path
            import pytesseract
            from PIL import Image
            
            images = convert_from_path(file_path)
            content = []
            for i, image in enumerate(images):
                text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                if text:
                    content.append(f"--- 第{i+1}页 ---\n{text}")
            if content:
                logger.info(f"PDF解析成功(OCR): {file_name}")
                return self._clean_content('\n'.join(content))
        except ImportError:
            logger.debug("pdf2image或pytesseract未安装")
        except Exception as e:
            logger.debug(f"PDF OCR失败: {e}")
        
        logger.warning(f"PDF解析失败: {file_name}")
        return f"[PDF文件: {file_name}]"
    
    def _parse_json(self, file_path: str) -> str:
        """解析JSON文件 - 直接读取文本"""
        return self._parse_txt(file_path)
    
    def _parse_csv(self, file_path: str) -> str:
        """解析CSV文件"""
        try:
            import csv
            content = []
            
            # 尝试不同编码
            encodings = ['utf-8', 'gbk', 'gb2312']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding, newline='') as f:
                        reader = csv.reader(f)
                        headers = next(reader, None)
                        if headers:
                            content.append(f"表头: {', '.join(headers)}")
                        
                        for i, row in enumerate(reader):
                            if i >= 100:  # 限制行数
                                content.append("... (更多数据)")
                                break
                            content.append(' | '.join(row))
                    break
                except UnicodeDecodeError:
                    continue
            
            if content:
                logger.info(f"CSV解析成功: {os.path.basename(file_path)}")
                return self._clean_content('\n'.join(content))
        except Exception as e:
            logger.warning(f"CSV解析失败: {os.path.basename(file_path)}")
        
        return ""
    
    def _parse_markdown(self, file_path: str) -> str:
        """解析Markdown文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 移除markdown标记但保留内容
            import re
            # 移除图片链接
            content = re.sub(r'!\[([^\]]*)\]\([^\)]+\)', r'[图片: \1]', content)
            # 移除普通链接
            content = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', content)
            # 移除格式标记
            content = re.sub(r'[#*_`~]', '', content)
            
            logger.info(f"Markdown解析成功: {os.path.basename(file_path)}")
            return self._clean_content(content)
        except Exception as e:
            logger.warning(f"Markdown解析失败: {os.path.basename(file_path)}")
            return ""
    
    def _parse_html(self, file_path: str) -> str:
        """解析HTML/XML文件"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
            
            # 移除script和style
            for script in soup(["script", "style"]):
                script.decompose()
            
            # 获取文本
            text = soup.get_text()
            
            # 清理
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            logger.info(f"HTML/XML解析成功: {os.path.basename(file_path)}")
            return self._clean_content(text)
        except ImportError:
            logger.debug("beautifulsoup4未安装")
            return self._parse_txt(file_path)
        except Exception as e:
            logger.warning(f"HTML解析失败: {os.path.basename(file_path)}")
            return ""
    
    def _parse_zip(self, file_path: str) -> str:
        """解析ZIP压缩包"""
        try:
            import zipfile
            content = []
            
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                file_list = zip_ref.namelist()
                content.append(f"压缩包内容 ({len(file_list)} 个文件):")
                
                for name in file_list[:20]:  # 只列前20个
                    content.append(f"  - {name}")
                if len(file_list) > 20:
                    content.append(f"  ... 还有 {len(file_list)-20} 个文件")
            
            logger.info(f"ZIP解析成功: {os.path.basename(file_path)}")
            return self._clean_content('\n'.join(content))
        except Exception as e:
            logger.warning(f"ZIP解析失败: {os.path.basename(file_path)}")
            return f"[ZIP文件: {os.path.basename(file_path)}]"
    
    def _parse_xmind(self, file_path: str) -> str:
        """解析XMind思维导图文件（支持旧版XML和新版JSON格式）"""
        try:
            import zipfile
            import json

            content = []

            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                names = zip_ref.namelist()

                # 优先尝试新版JSON格式（XMind 8+ / Zen）
                json_content = None
                for name in names:
                    if name == 'content.json' or name.endswith('/content.json'):
                        with zip_ref.open(name) as f:
                            json_content = json.loads(f.read().decode('utf-8'))
                        break

                if json_content:
                    topics = []
                    def extract_json_topics(node, level=0):
                        indent = '  ' * level
                        title = node.get('title', '')
                        if title:
                            topics.append(f"{indent}• {title}")
                        for child in node.get('children', {}).get('attached', []):
                            extract_json_topics(child, level + 1)
                        for child in node.get('children', {}).get('detached', []):
                            extract_json_topics(child, level + 1)

                    # content.json 是数组，每个元素是一个画布(sheet)
                    sheets = json_content if isinstance(json_content, list) else [json_content]
                    for sheet in sheets:
                        sheet_title = sheet.get('title', '')
                        if sheet_title:
                            content.append(f"=== {sheet_title} ===")
                        root_topic = sheet.get('rootTopic', {})
                        extract_json_topics(root_topic)

                    if topics:
                        content.extend(topics)
                        content.append(f"\n文件: {os.path.basename(file_path)}")
                        content.append(f"主题数: {len(topics)}")
                        logger.info(f"XMind解析成功(JSON): {os.path.basename(file_path)}")
                        return self._clean_content('\n'.join(content))

                # 旧版XML格式
                import xml.etree.ElementTree as ET
                content_xml = None
                for name in names:
                    if name == 'content.xml':
                        with zip_ref.open(name) as f:
                            content_xml = f.read().decode('utf-8')
                        break

                if content_xml:
                    root = ET.fromstring(content_xml)
                    ns = {'xmap': 'urn:xmind:xmap:xmlns:content:2.0'}

                    topics = []
                    def extract_xml_topics(topic_el, level=0):
                        title_el = topic_el.find('xmap:title', ns)
                        if title_el is not None and title_el.text:
                            topics.append(f"{'  ' * level}• {title_el.text}")
                        children_el = topic_el.find('xmap:children', ns)
                        if children_el is not None:
                            for topics_el in children_el.findall('xmap:topics', ns):
                                for child_topic in topics_el.findall('xmap:topic', ns):
                                    extract_xml_topics(child_topic, level + 1)

                    for sheet in root.findall('xmap:sheet', ns):
                        root_topic = sheet.find('xmap:topic', ns)
                        if root_topic is not None:
                            extract_xml_topics(root_topic)

                    if topics:
                        content.append("思维导图内容:")
                        content.extend(topics)
                        content.append(f"\n文件: {os.path.basename(file_path)}")
                        content.append(f"主题数: {len(topics)}")
                        logger.info(f"XMind解析成功(XML): {os.path.basename(file_path)}")
                        return self._clean_content('\n'.join(content))

                logger.warning(f"XMind文件无内容: {os.path.basename(file_path)}")
                return ""

        except Exception as e:
            logger.warning(f"XMind解析失败: {os.path.basename(file_path)}")
            return f"[XMind文件: {os.path.basename(file_path)}]"

    def _parse_pptx(self, file_path: str) -> str:
        """解析PPT演示文稿"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            content = []

            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                slide_texts.append(' | '.join(row_text))

                if slide_texts:
                    content.append(f"--- 第{i}页 ---")
                    content.extend(slide_texts)

            if content:
                logger.info(f"PPT解析成功: {os.path.basename(file_path)}")
                return self._clean_content('\n'.join(content))

            logger.warning(f"PPT无文本内容: {os.path.basename(file_path)}")
            return ""
        except ImportError:
            logger.warning("python-pptx未安装，无法解析PPT")
            return ""
        except Exception as e:
            logger.warning(f"PPT解析失败: {os.path.basename(file_path)}")
            return f"[PPT文件: {os.path.basename(file_path)}]"
    
    def _clean_content(self, content: str) -> str:
        """清理内容，去除无用字符，保留段落换行结构"""
        # 去除HTML标签
        content = re.sub(r'<[^>]+>', '', content)
        # 统一换行符
        content = content.replace('\r\n', '\n').replace('\r', '\n')
        # 保留换行符，只压缩同一行内的连续空白
        lines = content.split('\n')
        lines = [re.sub(r'[ \t]+', ' ', line).strip() for line in lines]
        content = '\n'.join(lines)
        # 连续空行最多保留一个（段落分隔）
        content = re.sub(r'\n{3,}', '\n\n', content)
        # 移除控制字符（保留 \n \t）
        content = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', content)

        return content.strip()
    
    def _extract_problem_info(self, file_name: str, content: str) -> Dict:
        """
        从文件名和内容中提取问题信息
        
        Args:
            file_name: 文件名
            content: 文件内容
            
        Returns:
            包含问题ID和问题类型的字典
        """
        info = {'problem_id': '', 'problem_type': ''}
        
        id_match = re.search(r'(\d{7,})', file_name)
        if id_match:
            info['problem_id'] = id_match.group(1)
        
        type_keywords = {
            '开户': '开户报错',
            '报错': '系统报错',
            '变更': '套餐变更',
            '移机': '移机改号',
            '改号': '移机改号',
            '销户': '销户问题',
            '销号': '销户问题',
            '合约': '合约问题',
            '解约': '合约问题',
            '费用': '费用问题',
            '扣费': '费用问题',
            '停机': '服务状态',
            '复机': '服务状态',
            '终端': '终端问题',
            '跨域': '跨域业务',
            '实名': '实名认证',
            '副卡': '副卡问题',
            '亲情': '亲情业务',
            '产品': '产品受理',
            '套餐': '套餐变更'
        }
        
        combined_text = file_name + ' ' + content[:500]
        for keyword, category in type_keywords.items():
            if keyword in combined_text:
                info['problem_type'] = category
                break
        
        if not info['problem_type']:
            info['problem_type'] = '其他问题'
        
        return info
    
    def parse_directory(self, dir_path: str, recursive: bool = True) -> List[Dict]:
        """
        解析目录下的所有文件
        
        Args:
            dir_path: 目录路径
            recursive: 是否递归解析子目录
            
        Returns:
            解析结果列表
        """
        results = []
        
        logger.info(f"开始解析目录: {dir_path}")
        
        if not os.path.exists(dir_path):
            logger.error(f"目录不存在: {dir_path}")
            return results
        
        if not os.path.isdir(dir_path):
            logger.error(f"路径不是目录: {dir_path}")
            return results
        
        try:
            if recursive:
                logger.info("使用递归模式解析")
                for root, dirs, files in os.walk(dir_path):
                    logger.info(f"处理目录: {root}, 文件数: {len(files)}")
                    for file in files:
                        try:
                            # 确保文件名正确编码
                            if isinstance(file, bytes):
                                file = file.decode('utf-8', errors='replace')
                            file_path = os.path.join(root, file)
                            logger.info(f"处理文件: {file_path}")
                            result = self.parse_file(file_path)
                            if result:
                                results.append(result)
                                logger.info(f"成功解析文件: {file}, 结果数: {len(results)}")
                        except Exception as e:
                            logger.error(f"跳过文件 {file}: {e}")
                            import traceback
                            traceback.print_exc()
            else:
                logger.info("使用非递归模式解析")
                for file in os.listdir(dir_path):
                    try:
                        if isinstance(file, bytes):
                            file = file.decode('utf-8', errors='replace')
                        file_path = os.path.join(dir_path, file)
                        if os.path.isfile(file_path):
                            logger.info(f"处理文件: {file_path}")
                            result = self.parse_file(file_path)
                            if result:
                                results.append(result)
                                logger.info(f"成功解析文件: {file}, 结果数: {len(results)}")
                    except Exception as e:
                        logger.error(f"跳过文件 {file}: {e}")
                        import traceback
                        traceback.print_exc()
        except Exception as e:
            logger.error(f"解析目录失败: {e}")
            import traceback
            traceback.print_exc()
        
        logger.info(f"共解析 {len(results)} 个文件")
        return results


def extract_solution(content: str) -> Tuple[str, str]:
    """
    从内容中提取问题描述和解决方案
    
    Args:
        content: 文件内容
        
    Returns:
        (问题描述, 解决方案) 元组
    """
    problem_desc = ""
    solution = ""
    
    solution_keywords = ['解决', '处理', '方案', '步骤', '方法', '操作', '建议', '修改', '配置']
    
    lines = content.split('\n')
    solution_start = -1
    
    for i, line in enumerate(lines):
        for keyword in solution_keywords:
            if keyword in line and len(line) < 50:
                solution_start = i
                break
        if solution_start >= 0:
            break
    
    if solution_start >= 0:
        problem_desc = '\n'.join(lines[:solution_start]).strip()
        solution = '\n'.join(lines[solution_start:]).strip()
    else:
        if len(content) > 500:
            problem_desc = content[:500]
            solution = content[500:]
        else:
            problem_desc = content
            solution = ""
    
    return problem_desc, solution


if __name__ == "__main__":
    parser = FileParser()
    
    test_dir = r"D:\工作文档\cbss2.0体验版\沃工单问题定位"
    results = parser.parse_directory(test_dir)
    
    print(f"解析文件数量: {len(results)}")
    if results:
        print(f"示例文件: {results[0]['file_name']}")
        print(f"问题类型: {results[0]['problem_type']}")
        print(f"内容预览: {results[0]['content'][:200]}...")
